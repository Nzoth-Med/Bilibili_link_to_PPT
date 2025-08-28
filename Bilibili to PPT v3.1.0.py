# -*- coding: utf-8 -*-
"""
Bilibili 链接一键生成PPT（控制台UI，无需 rich）
“少页高质量”版：显著增强去重，减少重复页

依赖：
    pip install yt-dlp opencv-python pillow imagehash python-pptx numpy
可选（更快的图像处理）：pip install opencv-contrib-python
"""

import os
import re
import sys
import time
import math
import threading
import shutil
from concurrent.futures import ThreadPoolExecutor, as_completed

import cv2
import numpy as np
import yt_dlp
from PIL import Image
import imagehash
from pptx import Presentation
from pptx.util import Inches
from shutil import which

# ========= 参数（少页高质量，更激进去重） =========
MAX_WORKERS_READ = max(4, os.cpu_count() or 4)
MAX_WORKERS_SAVE = max(4, (os.cpu_count() or 4) // 2)

# 目标抽样规模（少页模式更少）
TARGET_FRAMES_MIN = 40
TARGET_FRAMES_MAX = 90

# 场景最小时长（更长，减少近邻重复）
MIN_SCENE_SECONDS = 2.8

# 基线哈希阈值（后续会自适应微调）
BASE_PHASH_KEEP_DIFF = 18   # 相邻帧 pHash 差 >= 才认为有明显变化
BASE_DHASH_KEEP_DIFF = 18

# 二次清洗（滑窗）认为“相似”的阈值（更严格）
POST_PHASH_SIM = 12         # 小于此差值当作相似
POST_DHASH_SIM = 12
HIST_CORR_SIM = 0.94        # HSV 颜色相关系数 > 0.94 认为相似
HIST_CHISQR_DIFF = 0.20     # 卡方差异 < 0.20 认为相似（差异小）

TITLE_MAXLEN = 120
FRAME_JPEG_QUALITY = 92
IMG_MAX_WIDTH = 1920
TOTAL_BAR_WIDTH = 40
SLIDE_HEIGHT_INCH = 7.5     # 16:9 宽 13.33 高 7.5

# ========= ANSI 颜色与UI =========
C_RESET = "\033[0m"; C_BOLD = "\033[1m"
C_GREEN = "\033[92m"; C_CYAN = "\033[96m"; C_YELLOW = "\033[93m"
C_MAGENTA = "\033[95m"; C_BLUE = "\033[94m"; C_DIM = "\033[2m"; C_RED = "\033[91m"
BOX_TL, BOX_TR, BOX_BL, BOX_BR, BOX_H, BOX_V = "┏", "┓", "┗", "┛", "━", "┃"

def box_title(title: str):
    line = BOX_TL + BOX_H * (len(title) + 2) + BOX_TR
    print(C_CYAN + line + C_RESET)
    print(C_CYAN + BOX_V + C_RESET + " " + C_BOLD + title + C_RESET + " " + C_CYAN + BOX_V + C_RESET)
    print(C_CYAN + BOX_BL + BOX_H * (len(title) + 2) + BOX_BR + C_RESET)

def fmt_hms(seconds: float) -> str:
    seconds = int(round(seconds or 0))
    h = seconds // 3600; m = (seconds % 3600) // 60; s = seconds % 60
    parts = []
    if h > 0: parts.append(f"{h}h")
    if m > 0 or h > 0: parts.append(f"{m}min")
    parts.append(f"{s}s")
    return " ".join(parts)

def sanitize_filename(name: str) -> str:
    name = (name or "").strip()
    name = re.sub(r'[\\/:*?"<>|]', "·", name)
    name = re.sub(r'\s+', " ", name)
    name = name[:TITLE_MAXLEN].rstrip(". ")
    return name or "bilibili_video"

def progress_bar(prefix: str, ratio: float, width: int = 30) -> str:
    ratio = max(0.0, min(1.0, ratio))
    filled = int(width * ratio)
    bar = "█" * filled + "░" * (width - filled)
    return f"{prefix} [{bar}] {int(ratio*100):3d}%"

def print_refresh(line: str):
    cols = shutil.get_terminal_size(fallback=(100, 20)).columns
    sys.stdout.write("\r" + " " * (cols - 1))
    sys.stdout.write("\r" + line)
    sys.stdout.flush()

def print_step_header(step_name: str):
    print()
    print(C_MAGENTA + f"— {step_name} —" + C_RESET)

class TotalProgress:
    # 权重：下载0.35 抽帧+检测0.45 生成PPT0.15 清理0.05
    def __init__(self):
        self.weights = {"download": 0.35, "frames": 0.45, "ppt": 0.15, "cleanup": 0.05}
        self.progress = {k: 0.0 for k in self.weights}
        self.start = time.time()
    def update(self, key: str, ratio: float):
        self.progress[key] = max(0.0, min(1.0, ratio))
        total = sum(self.progress[k] * self.weights[k] for k in self.weights)
        print_refresh(progress_bar(C_BOLD + "总进度" + C_RESET, total, TOTAL_BAR_WIDTH) +
                      f"   已用时: {fmt_hms(time.time() - self.start)}")
        return total
TP = TotalProgress()

# ========= 下载&标题 =========
def extract_meta_and_title(url: str):
    """
    最小改动：在调用 yt_dlp.extract_info 时显示解析进度（spinner），
    返回值与原来完全相同： (sanitize_filename(title), duration)
    """
    ydl_opts = {"quiet": True, "no_warnings": True, "skip_download": True, "noprogress": True}

    result = {}
    exc = {}

    def worker():
        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(url, download=False)
            result["info"] = info
        except Exception as e:
            exc["error"] = e

    th = threading.Thread(target=worker, daemon=True)
    th.start()

    # 简单 spinner，用 print_refresh 原地刷新，不改动其它输出逻辑
    spinner = ["|", "/", "-", "\\"]
    i = 0
    t0 = time.time()
    try:
        while th.is_alive():
            elapsed = fmt_hms(time.time() - t0)
            print_refresh(C_CYAN + f"解析元信息 {spinner[i % len(spinner)]}  已用时: {elapsed}" + C_RESET)
            time.sleep(0.12)
            i += 1
    finally:
        # 解析完成后清理那一行（避免残留），不新增多余输出
        cols = shutil.get_terminal_size(fallback=(100, 20)).columns
        print_refresh(" " * (cols - 1))

    if exc:
        # 将后台异常抛回主线程，行为与直接调用 extract_info 保持一致（会被上层捕获）
        raise exc["error"]

    info = result.get("info", {}) or {}
    title = info.get("title") or "bilibili_video"
    duration = info.get("duration") or 0
    return sanitize_filename(title), duration

def ensure_folder(title: str):
    folder = os.path.abspath(title); os.makedirs(folder, exist_ok=True); return folder

def candidate_video_paths(folder: str, title: str):
    prefs = [".mp4", ".mkv", ".webm", ".flv", ".mov", ".avi"]
    paths = [os.path.join(folder, f"{title}{ext}") for ext in prefs]
    extra = []
    for f in os.listdir(folder):
        if f.startswith(title) and os.path.splitext(f)[1].lower() in prefs:
            extra.append(os.path.join(folder, f))
    seen, ordered = set(), []
    for p in paths + sorted(extra):
        if os.path.exists(p) and p not in seen:
            ordered.append(p); seen.add(p)
    return ordered

def is_openable_video(path: str) -> bool:
    cap = cv2.VideoCapture(path)
    ok = cap.isOpened()
    if ok:
        fps = cap.get(cv2.CAP_PROP_FPS) or 0
        frames = cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0
        ok = (fps > 0 and frames > 0)
    cap.release()
    return ok

def find_final_video_file(folder: str, title: str):
    # 1) 优先匹配标准后缀
    for p in candidate_video_paths(folder, title):
        if is_openable_video(p): return p
    # 2) 检测分轨，尝试合并
    video_part, audio_part = None, None
    for f in os.listdir(folder):
        low = f.lower()
        if f.startswith(title):
            if low.endswith(".m4s") or ("f" in low and low.endswith(".mp4")):
                video_part = os.path.join(folder, f)
            if low.endswith(".m4a"):
                audio_part = os.path.join(folder, f)
    if video_part and audio_part:
        if which("ffmpeg"):
            merged = os.path.join(folder, f"{title}.mp4")
            os.system(f'ffmpeg -y -i "{video_part}" -i "{audio_part}" -c copy "{merged}"')
            if os.path.exists(merged) and is_openable_video(merged): return merged
        raise RuntimeError("检测到分轨但系统缺少 ffmpeg 以合并，请安装并加入 PATH。")
    # 3) 兜底：挑可打开的最大文件
    files = [os.path.join(folder, f) for f in os.listdir(folder) if os.path.isfile(os.path.join(folder,f))]
    files.sort(key=lambda p: os.path.getsize(p), reverse=True)
    for p in files:
        if is_openable_video(p): return p
    return None

def download_video(url: str, folder: str, title: str):
    print_step_header("下载视频")
    start = time.time()
    # 防重复
    for p in candidate_video_paths(folder, title):
        if os.path.exists(p) and os.path.getsize(p) > 5 * 1024 * 1024 and is_openable_video(p):
            TP.update("download", 1.0)
            print_refresh(C_YELLOW + "已存在视频文件，跳过下载：" + p + C_RESET + " " * 10 + "\n")
            print(C_DIM + f"下载耗时: {fmt_hms(0)}  平均速度: 跳过" + C_RESET)
            return p

    def hook(d):
        if d['status'] == 'downloading':
            total = d.get('total_bytes') or d.get('total_bytes_estimate') or 0
            downloaded = d.get('downloaded_bytes', 0)
            speed = d.get('speed', 0.0) or 0.0
            eta = d.get('eta', None)
            ratio = (downloaded / total) if total else 0.0
            TP.update("download", min(0.99, ratio))
            if eta is None and speed > 0 and total:
                eta = max(0, int((total - downloaded) / speed))
            eta_str = fmt_hms(eta) if eta is not None else "--"
            spd_str = f"{downloaded/1024/1024:.1f}MB / {(total or 0)/1024/1024:.1f}MB @ {(speed or 0)/1024:.0f}KB/s"
            print_refresh(C_GREEN + progress_bar("下载", ratio, 32) + C_RESET +
                          f"  速度: {spd_str}  预计剩余: {eta_str}")

    ydl_opts = {
        "outtmpl": os.path.join(folder, f"{title}.%(ext)s"),
        "merge_output_format": "mp4",
        "format": "bv*+ba/b",
        "noprogress": True,
        "progress_hooks": [hook],
        "quiet": True,
        "no_warnings": True,
        "postprocessors": [{"key": "FFmpegVideoRemuxer", "preferedformat": "mp4"}],
    }
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        ydl.download([url])

    TP.update("download", 1.0)
    final_path = find_final_video_file(folder, title)
    if not final_path:
        raise RuntimeError("下载完成但未找到可打开的视频文件（可能缺少 ffmpeg 或下载失败）。")
    spent = time.time() - start
    size = os.path.getsize(final_path)
    avg = (size / spent / 1024) if spent > 0 else 0
    print("\n" + C_DIM + f"下载完成  耗时: {fmt_hms(spent)}  平均速度: {avg:.0f}KB/s" + C_RESET)
    return final_path

# ========= 视频分析 =========
def analyze_video(video_path: str):
    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        raise RuntimeError("无法打开视频：" + video_path)
    fps = cap.get(cv2.CAP_PROP_FPS) or 25.0
    frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
    width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH) or 0)
    height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT) or 0)
    duration = frames / fps if fps > 0 else 0
    cap.release()
    return fps, frames, duration, (width, height)

def auto_frame_interval(fps: float, frames: int, duration: float):
    if frames <= 0 or fps <= 0: return 10
    # 少页：抽样更稀疏
    target = int(max(TARGET_FRAMES_MIN, min(TARGET_FRAMES_MAX, max(1, duration // 3))))  # 约每3秒一帧
    interval = max(1, int(round(frames / max(1, target))))
    return interval

def read_frame_at(video_path: str, frame_idx: int):
    cap = cv2.VideoCapture(video_path)
    cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
    ok, frame = cap.read()
    cap.release()
    if not ok or frame is None: return None
    h, w = frame.shape[:2]
    if w > IMG_MAX_WIDTH:
        scale = IMG_MAX_WIDTH / w
        frame = cv2.resize(frame, (int(w*scale), int(h*scale)), interpolation=cv2.INTER_AREA)
    return frame

# ========= 图像度量 =========
def dhash_of_frame(frame_bgr: np.ndarray, hash_size: int = 8):
    # dHash：先转灰度，缩放到 (hash_size+1) x hash_size，比较相邻像素
    gray = cv2.cvtColor(frame_bgr, cv2.COLOR_BGR2GRAY)
    resized = cv2.resize(gray, (hash_size + 1, hash_size), interpolation=cv2.INTER_AREA)
    diff = resized[:, 1:] > resized[:, :-1]
    # 转换为 PIL 的 ImageHash 兼容对象
    return imagehash.ImageHash(diff)

def phash_of_frame(frame_bgr: np.ndarray):
    img = cv2.cvtColor(frame_bgr, cv2.COLOR_BGR2RGB)
    pil = Image.fromarray(img)
    return imagehash.phash(pil)

def hsv_histogram(frame_bgr: np.ndarray):
    hsv = cv2.cvtColor(frame_bgr, cv2.COLOR_BGR2HSV)
    hist = cv2.calcHist([hsv], [0,1,2], None, [32,8,8], [0,180, 0,256, 0,256])
    hist = cv2.normalize(hist, hist).flatten()
    return hist

def hist_compare_corr(h1, h2):
    # 相关系数（1.0 相同，越大越相似）
    return float(cv2.compareHist(h1.astype('float32'), h2.astype('float32'), cv2.HISTCMP_CORREL))

def hist_compare_chisqr(h1, h2):
    # 卡方（0 相同，越大越不同）
    return float(cv2.compareHist(h1.astype('float32'), h2.astype('float32'), cv2.HISTCMP_CHISQR))

def sharpness_score(frame_bgr: np.ndarray):
    gray = cv2.cvtColor(frame_bgr, cv2.COLOR_BGR2GRAY)
    lap = cv2.Laplacian(gray, cv2.CV_64F)
    var = lap.var()
    # 灰度熵（信息量）
    hist = cv2.calcHist([gray],[0],None,[256],[0,256]).ravel()
    hist = hist / (hist.sum() + 1e-8)
    entropy = -np.sum(hist * np.log2(hist + 1e-12))
    return 0.7 * var + 0.3 * (entropy * 1000.0)  # 同量纲混合

# ========= 核心：抽帧 + 变动检测 + 强力去重 =========
def ensure_empty_dir(p: str):
    if os.path.isdir(p): shutil.rmtree(p, ignore_errors=True)
    os.makedirs(p, exist_ok=True)

def extract_keyframes(video_path: str, fps: float, frames_total: int, interval: int, out_dir: str):
    print_step_header("抽帧与画面变动检测（少页高质量）")
    ensure_empty_dir(out_dir)

    indices = list(range(0, frames_total, interval))
    total = len(indices)
    t0 = time.time(); processed = 0; last_update = t0

    # 多线程读取 + 计算度量
    data = {}  # idx -> dict(frame, phash, dhash, hist, sharp)
    def worker(idx):
        st = time.time()
        f = read_frame_at(video_path, idx)
        if f is None: return idx, None, time.time() - st
        ph = phash_of_frame(f)
        dh = dhash_of_frame(f)
        hs = hsv_histogram(f)
        sp = sharpness_score(f)
        return idx, (f, ph, dh, hs, sp), time.time() - st

    with ThreadPoolExecutor(max_workers=MAX_WORKERS_READ) as ex:
        futures = [ex.submit(worker, i) for i in indices]
        durations = []
        for fut in as_completed(futures):
            idx, packed, dt = fut.result()
            if packed is not None:
                data[idx] = {
                    "frame": packed[0], "phash": packed[1], "dhash": packed[2],
                    "hist": packed[3], "sharp": packed[4]
                }
            processed += 1; durations.append(dt)
            now = time.time()
            if now - last_update >= 0.15 or processed == total:
                ratio = processed / total
                TP.update("frames", 0.10 * ratio)
                elapsed = now - t0
                fps_cur = processed / elapsed if elapsed > 0 else 0
                eta = (total - processed) / fps_cur if fps_cur > 0 else 0
                print_refresh(C_GREEN + progress_bar("抽帧+度量", ratio, 32) + C_RESET +
                              f"  速度: {fps_cur:.1f}帧/s  预计剩余: {fmt_hms(eta)}")
                last_update = now

    order = sorted(data.keys())
    if not order: return []

    # 计算相邻差异（pHash 为主）以便自适应阈值
    ph_diffs = []
    for i in range(1, len(order)):
        ph_diffs.append(data[order[i]]["phash"] - data[order[i-1]]["phash"])
    median_diff = int(np.median(ph_diffs)) if ph_diffs else BASE_PHASH_KEEP_DIFF
    # 自适应：但保持偏“严格”
    PHASH_KEEP_DIFF = max(BASE_PHASH_KEEP_DIFF, int(0.8 * median_diff + 8))

    # —— 第一阶段：分段（根据哈希 + 最小场景间隔），每段挑最清晰的一帧 ——
    segments = []
    cur_seg = [order[0]]
    last_cut_idx = order[0]
    min_scene_frames = int(MIN_SCENE_SECONDS * fps)
    for i in range(1, len(order)):
        a, b = order[i-1], order[i]
        ph = data[b]["phash"] - data[a]["phash"]
        dh = data[b]["dhash"] - data[a]["dhash"]
        # 颜色差异
        corr = hist_compare_corr(data[a]["hist"], data[b]["hist"])  # 大=相似
        chi = hist_compare_chisqr(data[a]["hist"], data[b]["hist"]) # 小=相似
        strong_change = (ph >= PHASH_KEEP_DIFF) or (dh >= BASE_DHASH_KEEP_DIFF) or (corr < 0.85) or (chi > 0.35)
        long_enough = (b - last_cut_idx) >= min_scene_frames
        if strong_change or long_enough:
            segments.append(cur_seg)
            cur_seg = [b]
            last_cut_idx = b
        else:
            cur_seg.append(b)
    if cur_seg: segments.append(cur_seg)

    # 在每个 segment 中选最佳（清晰度+信息量）
    candidates = []
    for seg in segments:
        if not seg: continue
        best_idx = max(seg, key=lambda k: data[k]["sharp"])
        candidates.append(best_idx)

    # —— 第二阶段：滑窗相似性去重（更激进），避免近邻重复 ——
    keep = []
    WINDOW = 5  # 与最近 5 个候选比较
    for idx in sorted(candidates):
        similar = False
        for prev in keep[-WINDOW:]:
            ph = data[idx]["phash"] - data[prev]["phash"]
            dh = data[idx]["dhash"] - data[prev]["dhash"]
            corr = hist_compare_corr(data[idx]["hist"], data[prev]["hist"])
            chi = hist_compare_chisqr(data[idx]["hist"], data[prev]["hist"])
            if (ph < POST_PHASH_SIM and dh < POST_DHASH_SIM and corr > HIST_CORR_SIM and chi < HIST_CHISQR_DIFF):
                similar = True; break
        if not similar:
            keep.append(idx)

    # 保存关键帧（多线程）
    print()
    print(C_YELLOW + f"自适应阈值：PHASH_KEEP_DIFF={PHASH_KEEP_DIFF}  最小场景时长={MIN_SCENE_SECONDS:.1f}s" + C_RESET)
    print(C_YELLOW + f"候选段数: {len(segments)}  候选帧: {len(candidates)}  二次去重后保留: {len(keep)}" + C_RESET)

    saved_paths = []
    save_start = time.time(); save_count = 0

    def save_worker(i, idx):
        frame = data[idx]["frame"]
        outp = os.path.join(out_dir, f"frame_{i:05d}.jpg")
        pil = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
        pil.save(outp, format="JPEG", quality=FRAME_JPEG_QUALITY, optimize=True)
        return outp

    with ThreadPoolExecutor(max_workers=MAX_WORKERS_SAVE) as ex:
        futures = {ex.submit(save_worker, i, idx): (i, idx) for i, idx in enumerate(keep)}
        total_save = len(futures) if futures else 1
        done_last = 0; t_last = save_start
        for fut in as_completed(futures):
            path = fut.result()
            saved_paths.append(path); save_count += 1
            ratio = save_count / max(1, total_save)
            TP.update("frames", 0.10 + 0.90 * ratio)
            now = time.time()
            elapsed = now - save_start
            spd = save_count / elapsed if elapsed > 0 else 0
            eta = (total_save - save_count) / max(spd, 1e-6)
            print_refresh(C_GREEN + progress_bar("去重&保存关键帧", ratio, 32) + C_RESET +
                          f"  速度: {spd:.1f}页/s  预计剩余: {fmt_hms(eta)}")
            done_last = save_count; t_last = now

    total_spent = time.time() - t0
    print("\n" + C_DIM + f"抽帧+检测+强力去重完成  耗时: {fmt_hms(total_spent)}  "
          f"平均处理速度: {(len(indices)/total_spent if total_spent>0 else 0):.1f}帧/s  "
          f"最终关键帧: {len(saved_paths)}" + C_RESET)

    saved_paths.sort()
    return saved_paths

# ========= 生成 PPT =========
def build_ppt_from_images(images: list, out_ppt_path: str, title: str):
    print_step_header("生成PPT")
    t0 = time.time()
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCH)

    # 标题页
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "由脚本自动生成"

    total = len(images); done = 0; last = t0
    for img in images:
        blank = prs.slide_layouts[6]
        s = prs.slides.add_slide(blank)

        # === 修改：图片直接铺满幻灯片 ===
        left, top = Inches(0), Inches(0)
        pic = s.shapes.add_picture(
            img,
            left,
            top,
            width=prs.slide_width,
            height=prs.slide_height
        )

        done += 1
        TP.update("ppt", done / max(1, total))
        now = time.time()
        if now - last >= 0.1 or done == total:
            elapsed = now - t0; spd = done / elapsed if elapsed > 0 else 0
            eta = (total - done) / max(spd, 1e-6)
            print_refresh(C_GREEN + progress_bar("写入PPT", done / max(1, total), 32) + C_RESET +
                          f"  速度: {spd:.1f}页/s  预计剩余: {fmt_hms(eta)}")
            last = now

    prs.save(out_ppt_path)
    spent = time.time() - t0
    avg = (total / spent) if spent > 0 else 0
    print("\n" + C_DIM + f"PPT生成完成  耗时: {fmt_hms(spent)}  平均速度: {avg:.1f}页/s" + C_RESET)
    print(C_BOLD + f"PPT页数: {len(prs.slides)}（含标题页）" + C_RESET)
    return len(prs.slides)


# ========= 清理 =========
def cleanup_temp(temp_dir: str):
    print_step_header("清理临时文件")
    t0 = time.time()
    if os.path.isdir(temp_dir):
        files = [os.path.join(temp_dir, f) for f in os.listdir(temp_dir)]
        total = max(1, len(files)); done = 0
        for p in files:
            try: os.remove(p)
            except: pass
            done += 1
            TP.update("cleanup", done / total)
            print_refresh(C_GREEN + progress_bar("清理", done / total, 32) + C_RESET + " 正在删除临时文件...")
        shutil.rmtree(temp_dir, ignore_errors=True)
    TP.update("cleanup", 1.0)
    print("\n" + C_DIM + f"清理完成  耗时: {fmt_hms(time.time() - t0)}" + C_RESET)

# ========= 主流程 =========
def main():
    os.system("")  # 使 ANSI 颜色在 Windows 上生效（启用颜色）
    box_title("Bilibili 视频一键生成PPT（少页高质量）")

    try:
        url = input(C_BOLD + "请输入B站视频链接: " + C_RESET).strip().strip('"').strip("'")
        if not url:
            print(C_RED + "未输入链接，退出。" + C_RESET)
            input(C_BOLD + "\n按回车键退出..." + C_RESET); return

        # ========== 新增（最小改动）：先检查是否为合集（playlist） ==========
        print_step_header("解析视频信息")
        # 我们用 yt_dlp 抽取 info 来判断是否存在 entries（合集）
        ydl_opts = {"quiet": True, "no_warnings": True, "skip_download": True, "noprogress": True}
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info_all = ydl.extract_info(url, download=False)

        process_list = []  # 列表元素为 (entry_url, p_index_or_None)
        if info_all and isinstance(info_all, dict) and info_all.get("entries"):
            entries = [e for e in info_all.get("entries") if e]
            if entries:
                print(C_YELLOW + f"检测到合集，共 {len(entries)} 个分P" + C_RESET)
                for idx, entry in enumerate(entries, start=1):
                    # 优先使用 webpage_url / url 字段
                    entry_url = entry.get("webpage_url") or entry.get("url")
                    if not entry_url:
                        eid = entry.get("id")
                        if eid:
                            entry_url = f"https://www.bilibili.com/video/{eid}"
                        else:
                            entry_url = url  # 兜底回退到原链接（非常罕见）
                    process_list.append((entry_url, idx))
        else:
            process_list.append((url, None))

        # ========== 对 process_list 中的每一个项，逐个走原本的流程（未改动其他逻辑） ==========
        for entry_url, p_idx in process_list:
            meta_t0 = time.time()
            title, duration_meta = extract_meta_and_title(entry_url)
            # 如果是合集分P，给标题添加后缀以区分文件夹与输出
            if p_idx is not None:
                title = f"{title}_P{p_idx}"
            folder = ensure_folder(title)
            temp_dir = os.path.join(folder, "_tmp_frames")
            print(C_BOLD + f"视频标题: {title}" + C_RESET)
            print(C_DIM + f"输出目录: {folder}" + C_RESET)
            print(C_DIM + f"时长(来自元数据): {fmt_hms(duration_meta)}" + C_RESET)
            if not which("ffmpeg"):
                print(C_YELLOW + "提示：未检测到 ffmpeg，可正常下载，但如遇分轨将无法自动合并。建议安装 ffmpeg 并添加到 PATH。" + C_RESET)

            # 下载（防重复）
            video_path = download_video(entry_url, folder, title)

            # 分析视频
            fps, frames_total, duration_real, (w, h) = analyze_video(video_path)
            interval = auto_frame_interval(fps, frames_total, duration_real)
            print(C_DIM + f"视频参数: {w}x{h} @ {fps:.3f}fps, 总帧数≈{frames_total}, 实测时长: {fmt_hms(duration_real)}" + C_RESET)
            print(C_DIM + f"自动计算 frame_interval: {interval} （约每 {interval/max(fps,1e-6):.2f}s 一帧）" + C_RESET)

            # 抽帧 + 画面变动检测 + 强力去重 + 保存关键帧
            images = extract_keyframes(video_path, fps, frames_total, interval, temp_dir)

            # 生成PPT（每页仅一张图）
            ppt_path = os.path.join(folder, f"{title}.pptx")
            TP.update("ppt", 0.0)
            pages = build_ppt_from_images(images, ppt_path, title)
            print(C_GREEN + C_BOLD + f"已生成PPT：{ppt_path}" + C_RESET)

            # 清理
            cleanup_temp(temp_dir)

            # 汇总
            print("\n" + C_BLUE + "— 任务总结 —" + C_RESET)
            print(C_BOLD + f"标题: {title}" + C_RESET)
            print(C_BOLD + f"视频: {video_path}" + C_RESET)
            print(C_BOLD + f"PPT:   {ppt_path}" + C_RESET)
            print(C_BOLD + f"PPT页数: {pages}" + C_RESET)
            print(C_DIM + f"解析信息耗时: {fmt_hms(time.time() - meta_t0)}" + C_RESET)
                    # 汇总
            print("\n" + C_BLUE + "— 任务总结 —" + C_RESET)
            print(C_BOLD + f"标题: {title}" + C_RESET)
            print(C_BOLD + f"视频: {video_path}" + C_RESET)
            print(C_BOLD + f"PPT:   {ppt_path}" + C_RESET)
            print(C_BOLD + f"PPT页数: {pages}" + C_RESET)
            print(C_DIM + f"解析信息耗时: {fmt_hms(time.time() - meta_t0)}" + C_RESET)

            # === 新增：写入运行信息到 txt 文件 ===
            txt_path = os.path.join(folder, "脚本运行信息.txt")
            with open(txt_path, "w", encoding="utf-8") as f:
                f.write("— 脚本运行信息 —\n")
                f.write(f"标题: {title}\n")
                f.write(f"视频文件: {video_path}\n")
                f.write(f"PPT文件: {ppt_path}\n")
                f.write(f"PPT页数: {pages}\n")
                f.write(f"解析信息耗时: {fmt_hms(time.time() - meta_t0)}\n")
                f.write(f"总耗时: {fmt_hms(time.time() - TP.start)}\n")

            print(C_GREEN + C_BOLD + f"运行信息已保存到: {txt_path}" + C_RESET)


        # 在处理完所有分P 或 单视频后，再做统一的退出等待（保持原来交互）
        input(C_BOLD + "\n全部完成。按回车键退出..." + C_RESET)

    except Exception as e:
        print("\n" + C_RED + "发生错误: " + str(e) + C_RESET)
        input(C_BOLD + "\n按回车键退出..." + C_RESET)

if __name__ == "__main__":
    main()
